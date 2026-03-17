"""Kopiert Folien zwischen PowerPoint-Präsentationen.

Baut die Ziel-Präsentation auf, indem Folien aus verschiedenen Quellen
in der richtigen Reihenfolge kopiert werden.
"""

import copy
import logging
import os

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.package import PackURI

log = logging.getLogger(__name__)

# XML Namespaces
nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# Globaler Zähler für eindeutige Mediennamen
_media_counter = 0


def _find_layout(target_prs, layout_name, master_index=None):
    """Findet ein Layout in der Ziel-Präsentation.

    Args:
        target_prs: Ziel-Presentation
        layout_name: Name des gesuchten Layouts
        master_index: Optional - nur in diesem Master suchen
    """
    if master_index is not None:
        master = target_prs.slide_masters[master_index]
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout

    # Suche in allen Masters
    for master in target_prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout

    # Fallback
    return target_prs.slide_layouts[0]


def _find_song_master_index(target_prs):
    """Findet den Song-Master im Template (pattFill mit pct5 Hintergrund)."""
    for mi, master in enumerate(target_prs.slide_masters):
        bg = master._element.find(
            f'{{{nsmap["p"]}}}cSld/{{{nsmap["p"]}}}bg')
        if bg is not None:
            patt = bg.find('.//{%s}pattFill' % nsmap['a'])
            if patt is not None and patt.get('prst') == 'pct5':
                return mi
    return None


def copy_slide(target_prs, source_slide, source_prs, is_song=False,
               song_master_idx=None):
    """Kopiert eine einzelne Folie von der Quelle in die Ziel-Präsentation.

    Args:
        target_prs: Ziel-Presentation-Objekt
        source_slide: Quell-Slide-Objekt
        source_prs: Quell-Presentation-Objekt
        is_song: True wenn die Quelle eine Song-Datei ist
        song_master_idx: Master-Index für Song-Folien im Template
    """
    source_layout_name = source_slide.slide_layout.name

    if is_song and song_master_idx is not None:
        # Song-Folien: Layout vom Song-Master im Template verwenden
        layout = _find_layout(target_prs, source_layout_name, song_master_idx)
    else:
        # Template-Folien: Gleichnamiges Layout finden
        layout = _find_layout(target_prs, source_layout_name)

    new_slide = target_prs.slides.add_slide(layout)

    # Placeholder vom Layout entfernen (wir kopieren eigene Shapes)
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Hintergrund kopieren (falls die Folie einen eigenen hat)
    _copy_slide_background(new_slide, source_slide)

    # Alle Shapes kopieren
    for shape in source_slide.shapes:
        _copy_shape(new_slide, shape, source_slide, is_song=is_song)

    return new_slide


def _copy_slide_background(new_slide, source_slide):
    """Kopiert den Folienhintergrund (nur wenn die Quell-Folie einen eigenen hat)."""
    src_bg = source_slide._element.find(
        f'{{{nsmap["p"]}}}cSld/{{{nsmap["p"]}}}bg')
    if src_bg is None:
        return

    new_bg = copy.deepcopy(src_bg)

    # Bild-Referenzen im Hintergrund aktualisieren
    bgFill = new_bg.find('.//a:blipFill/a:blip', nsmap)
    if bgFill is not None:
        rId = bgFill.get('{%s}embed' % nsmap['r'])
        if rId:
            new_rId = _copy_image_rel(source_slide, new_slide, rId)
            if new_rId:
                bgFill.set('{%s}embed' % nsmap['r'], new_rId)

    cSld = new_slide._element.find('p:cSld', nsmap)
    if cSld is not None:
        existing_bg = cSld.find('p:bg', nsmap)
        if existing_bg is not None:
            cSld.remove(existing_bg)
        cSld.insert(0, new_bg)


def _copy_shape(new_slide, shape, source_slide, is_song=False):
    """Kopiert ein Shape inkl. Placeholder-Geometrie und Bild-Relationen."""
    new_el = copy.deepcopy(shape._element)

    if not is_song:
        # Placeholder-Geometrie explizit setzen (Position/Größe)
        # damit Shapes nicht von einem falschen Layout erben
        _resolve_placeholder_geometry(new_el, shape)

    # Bild-Relationen aktualisieren
    blips = new_el.findall('.//{%s}blip' % nsmap['a'])
    for blip in blips:
        rId = blip.get('{%s}embed' % nsmap['r'])
        if rId:
            new_rId = _copy_image_rel(source_slide, new_slide, rId)
            if new_rId:
                blip.set('{%s}embed' % nsmap['r'], new_rId)

    new_slide.shapes._spTree.append(new_el)


def _resolve_placeholder_geometry(new_el, shape):
    """Setzt explizite Position/Größe für Placeholder-Shapes.

    Entfernt die Placeholder-Referenz damit die Position nicht vom
    Ziel-Layout überschrieben wird.
    """
    ph = new_el.find('.//p:nvSpPr/p:nvPr/p:ph', nsmap)
    if ph is None:
        return

    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    if left is None or top is None or width is None or height is None:
        return

    spPr = new_el.find('.//p:spPr', nsmap)
    if spPr is None:
        return

    xfrm = spPr.find('a:xfrm', nsmap)
    if xfrm is None:
        xfrm = etree.SubElement(spPr, '{%s}xfrm' % nsmap['a'])
        spPr.insert(0, xfrm)

    off = xfrm.find('a:off', nsmap)
    if off is None:
        off = etree.SubElement(xfrm, '{%s}off' % nsmap['a'])
    off.set('x', str(left))
    off.set('y', str(top))

    ext = xfrm.find('a:ext', nsmap)
    if ext is None:
        ext = etree.SubElement(xfrm, '{%s}ext' % nsmap['a'])
    ext.set('cx', str(width))
    ext.set('cy', str(height))

    # Placeholder-Referenz entfernen
    ph.getparent().remove(ph)


def _copy_image_rel(source_slide, target_slide, source_rId):
    """Kopiert eine Bild-Relation von Quell- zu Ziel-Slide."""
    global _media_counter
    try:
        src_rel = source_slide.part.rels[source_rId]
        image_blob = src_rel.target_part.blob
        content_type = src_rel.target_part.content_type

        _media_counter += 1
        ext = _content_type_to_ext(content_type)
        media_name = f"image_auto_{_media_counter}{ext}"
        partname = PackURI(f'/ppt/media/{media_name}')

        image_part = Part(partname, content_type, target_slide.part.package, image_blob)
        new_rId = target_slide.part.relate_to(image_part, RT.IMAGE)
        return new_rId
    except (KeyError, AttributeError) as e:
        log.debug(f"Bild-Relation konnte nicht kopiert werden: {e}")
        return None


def _content_type_to_ext(content_type: str) -> str:
    """Konvertiert Content-Type zu Dateiendung."""
    mapping = {
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/tiff': '.tiff',
        'image/svg+xml': '.svg',
        'image/x-emf': '.emf',
        'image/x-wmf': '.wmf',
    }
    return mapping.get(content_type, '.png')


def build_presentation_from_plan(template_path: str, slide_plan: list,
                                 output_path: str) -> Presentation:
    """Baut eine Präsentation gemäß dem Folienplan."""
    template_prs = Presentation(template_path)

    import shutil
    shutil.copy2(template_path, output_path)
    target_prs = Presentation(output_path)

    # Song-Master im Template finden
    song_master_idx = _find_song_master_index(target_prs)
    if song_master_idx is not None:
        log.info(f"Song-Master gefunden: Master {song_master_idx}")
    else:
        log.warning("Kein Song-Master im Template gefunden!")

    # Alle bestehenden Folien entfernen
    while len(target_prs.slides) > 0:
        _delete_first_slide(target_prs)

    # Song-Dateien cachen
    song_cache = {}

    template_indices = set()

    for entry in slide_plan:
        if entry[0] == "template":
            slide_idx = entry[1]
            source_slide = template_prs.slides[slide_idx]
            template_indices.add(len(target_prs.slides))
            copy_slide(target_prs, source_slide, template_prs,
                       is_song=False)

        elif entry[0] == "file":
            file_path = entry[1]
            if file_path not in song_cache:
                song_cache[file_path] = Presentation(file_path)
            song_prs = song_cache[file_path]
            for slide in song_prs.slides:
                copy_slide(target_prs, slide, song_prs,
                           is_song=True, song_master_idx=song_master_idx)

        elif entry[0] == "skip":
            continue

    log.info(f"Präsentation erstellt: {len(target_prs.slides)} Folien")
    target_prs.save(output_path)
    return target_prs, template_indices


def _delete_first_slide(prs):
    """Löscht die erste Folie der Präsentation."""
    sldIdLst = prs.slides._sldIdLst
    if len(sldIdLst) == 0:
        return

    sldId = sldIdLst[0]
    rId = sldId.get('{%s}id' % nsmap['r'])
    sldIdLst.remove(sldId)

    if rId:
        try:
            prs.part.drop_rel(rId)
        except (KeyError, AttributeError):
            pass
