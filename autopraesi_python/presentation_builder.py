"""Baut die Gottesdienst-Präsentation zusammen."""

import logging
import os
import shutil

import copy

from lxml import etree
from pptx import Presentation

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from config import (TEMPLATE_PATH, OUTPUT_DIR_DESKTOP, OUTPUT_DIR_DROPBOX,
                    DEFAULT_SECTION_ORDER, SECTION_BLOCKS, INTRO_SLIDES,
                    TOGGLEABLE_SECTIONS)
from slide_copier import build_presentation_from_plan, copy_slide
from bible_fetcher import fetch_bible_text

# XML Namespaces
_nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

log = logging.getLogger(__name__)

# Template-Folien-Positionen für Lieder (0-basiert)
# Jede gibt an, welcher Template-Slide-Index durch ein Lied ersetzt wird
SONG_SLOTS = {
    "song1": 4,     # Template Folie 5 (0-basiert: 4)
    "song2": 9,     # Template Folie 10
    "song3": 14,    # Template Folie 15
    "song4": 21,    # Template Folie 22
    "song5": 25,    # Template Folie 26
    "song6": 31,    # Template Folie 32
    "song7": 35,    # Template Folie 36
}


def _replace_text_in_shape(shape, old_text, new_text):
    """Ersetzt Text in einem Shape, behält Formatierung bei."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for para in shape.text_frame.paragraphs:
        full_para_text = para.text
        if old_text in full_para_text:
            if len(para.runs) >= 1:
                combined = full_para_text.replace(old_text, new_text)
                para.runs[0].text = combined
                for run in para.runs[1:]:
                    run.text = ""
                replaced = True
    return replaced


def _replace_text_in_slide(slide, old_text, new_text):
    """Ersetzt Text in allen Shapes einer Folie."""
    for shape in slide.shapes:
        _replace_text_in_shape(shape, old_text, new_text)


def build_presentation(data, song_paths: dict, image_path: str = None,
                       fetch_bible: bool = True, output_name: str = None,
                       skip_slides: set = None, section_order: list = None,
                       extra_songs: dict = None, text_color: str = "white",
                       title_layout: dict = None, subtitle_layout: dict = None) -> str:
    """Erstellt die komplette Gottesdienst-Präsentation.

    Args:
        data: GodiPlanData-Objekt
        song_paths: Dict {slot_key: pfad} für die Lieder (inkl. Extras)
        image_path: Pfad zum Hintergrundbild (optional)
        fetch_bible: Ob Bibeltexte geladen werden sollen
        output_name: Dateiname (optional, z.B. 'Pa 18.3_ungeprüft.pptx')
        skip_slides: Set von 0-basierten Template-Folien-Indizes die übersprungen werden
        section_order: Liste von Section-Keys in gewünschter Reihenfolge
        extra_songs: Dict {slot_key: pfad} für Extra-Lieder (song_extra1, ...)

    Returns:
        Pfad zur erstellten Präsentation
    """
    # Ausgabe-Dateiname
    if output_name is None:
        day = data.date_str.split(".")[0] if data.date_str else "XX"
        month = data.date_str.split(".")[1] if data.date_str else "XX"
        output_name = f"So {day.lstrip('0')}.{month.lstrip('0')}_ungeprüft.pptx"
    output_path = os.path.join(OUTPUT_DIR_DESKTOP, output_name)

    # Extra-Songs: song_extra* Pfade aus song_paths extrahieren
    all_extra = {k: v for k, v in song_paths.items() if k.startswith("song_extra")}
    if extra_songs:
        all_extra.update(extra_songs)
    regular_paths = {k: v for k, v in song_paths.items() if not k.startswith("song_extra")}

    # === Phase 1: Folienplan erstellen ===
    slide_plan = _build_slide_plan(regular_paths, skip_slides=skip_slides,
                                   section_order=section_order,
                                   extra_song_paths=all_extra)

    # === Phase 2: Präsentation aus Plan bauen ===
    log.info("Baue Präsentation aus Folienplan...")
    prs, template_indices = build_presentation_from_plan(TEMPLATE_PATH, slide_plan, output_path)

    # === Phase 3: Hintergrundbild auf Thema-Folien setzen ===
    if image_path and os.path.exists(image_path):
        _set_theme_image(prs, image_path)

    # === Phase 4: Texte ersetzen ===
    _fill_all_text(prs, data, fetch_bible, template_indices)

    # === Phase 4b: Textfarbe auf Thema-Folien anpassen ===
    if text_color == "black":
        _set_theme_text_color(prs, "000000")

    # === Phase 4c: Text-Layout auf Thema-Folien anpassen ===
    if title_layout or subtitle_layout:
        _apply_theme_text_layout(prs, title_layout, subtitle_layout)

    # === Phase 5: Speichern ===
    prs.save(output_path)
    log.info(f"Präsentation gespeichert: {output_path} ({len(prs.slides)} Folien)")

    # Kopie in Dropbox ablegen
    dropbox_path = os.path.join(OUTPUT_DIR_DROPBOX, output_name)
    try:
        shutil.copy2(output_path, dropbox_path)
        log.info(f"Kopie in Dropbox: {dropbox_path}")
    except Exception as e:
        log.warning(f"Dropbox-Kopie fehlgeschlagen: {e}")

    return output_path


def _set_theme_image(prs, image_path: str):
    """Ersetzt das Template-Hintergrundbild auf allen Folien.

    Findet das ursprüngliche Template-Bild (224KB Platzhalter) und
    ersetzt es überall durch das neue Sonntagsbild.
    """
    with open(image_path, 'rb') as f:
        new_blob = f.read()

    # Alle Image-Parts durchgehen und das Template-Bild ersetzen
    replaced_parts = set()
    count = 0

    for slide in prs.slides:
        blips = slide._element.findall('.//{%s}blip' % _nsmap['a'])
        for blip in blips:
            rId = blip.get('{%s}embed' % _nsmap['r'])
            if not rId:
                continue
            try:
                rel = slide.part.rels[rId]
                part_id = id(rel.target_part)
                if part_id in replaced_parts:
                    count += 1
                    continue
                # Nur das Template-Bild ersetzen (nicht Song-Logos etc.)
                # Template-Bild ist deutlich kleiner als das Sonntagsbild
                if len(rel.target_part._blob) < len(new_blob):
                    rel.target_part._blob = new_blob
                    replaced_parts.add(part_id)
                    count += 1
            except (KeyError, AttributeError) as e:
                log.warning(f"Bild konnte nicht ersetzt werden: {e}")

    log.info(f"Hintergrundbild auf {count} Folien ersetzt")


def _set_theme_text_color(prs, hex_color: str):
    """Ändert die Textfarbe auf allen Thema-Folien (mit Hintergrundbild).

    Findet Folien die das Thema/Datum enthalten und setzt die Textfarbe
    aller Runs auf den angegebenen Hex-Wert.
    """
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    count = 0

    for slide in prs.slides:
        slide_text = _get_slide_text(slide)
        # Thema-Folien erkennen: enthalten "Gottesdienst am" und Theme-Text
        if "Gottesdienst am" not in slide_text:
            continue

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(f'{{{a_ns}}}rPr')
                    if rPr is None:
                        rPr = etree.SubElement(run._r, f'{{{a_ns}}}rPr')
                        run._r.insert(0, rPr)

                    # Bestehende Farbdefinition entfernen
                    for old_fill in rPr.findall(f'{{{a_ns}}}solidFill'):
                        rPr.remove(old_fill)

                    # Neue Farbe setzen
                    solidFill = etree.SubElement(rPr, f'{{{a_ns}}}solidFill')
                    srgbClr = etree.SubElement(solidFill, f'{{{a_ns}}}srgbClr')
                    srgbClr.set('val', hex_color)

        count += 1

    log.info(f"Textfarbe auf {count} Thema-Folien auf #{hex_color} gesetzt")


def _apply_theme_text_layout(prs, title_layout: dict = None, subtitle_layout: dict = None):
    """Passt Position und Größe der Textboxen auf Thema-Folien an.

    Layout-Werte sind in % der Foliengröße.
    fontSize ist in cqi-Einheiten (Container Query Inline) im Frontend,
    hier umgerechnet in PowerPoint-Punkte.
    """
    from pptx.util import Emu, Pt

    slide_w = prs.slide_width  # 9144000 EMU
    slide_h = prs.slide_height  # 6858000 EMU
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # fontSize: Frontend cqi → pt Umrechnung
    # 6.6 cqi entspricht 66pt bei voller Folie
    CQI_TO_PT = 10.0  # 1 cqi ≈ 10pt

    for slide in prs.slides:
        slide_text = _get_slide_text(slide)
        if "Gottesdienst am" not in slide_text:
            continue

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_text = shape.text_frame.text.strip()
            layout = None

            # Titel erkennen: das Shape mit dem Theme-Text (groß, oben)
            if title_layout and shape.name == "Titel 1":
                layout = title_layout
            elif subtitle_layout and shape.name == "CustomShape 1":
                layout = subtitle_layout

            # Fallback: nach Position identifizieren
            if not layout and title_layout:
                if shape.top < slide_h * 0.4 and shape.height > slide_h * 0.15:
                    layout = title_layout
            if not layout and subtitle_layout:
                if shape.top > slide_h * 0.6:
                    layout = subtitle_layout

            if not layout:
                continue

            # Position und Größe setzen
            shape.left = int(slide_w * layout["x"] / 100)
            shape.top = int(slide_h * layout["y"] / 100)
            shape.width = int(slide_w * layout["w"] / 100)
            shape.height = int(slide_h * layout["h"] / 100)

            # Schriftgröße
            if "fontSize" in layout:
                pt_size = int(layout["fontSize"] * CQI_TO_PT)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(pt_size)

            # Text horizontal und vertikal zentrieren
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            for para in shape.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
            # Vertikale Zentrierung
            txBody = shape._element.find(f'{{{a_ns}}}txBody')
            if txBody is None:
                txBody = shape._element.find('.//{%s}txBody' % a_ns)
            if txBody is not None:
                bodyPr = txBody.find(f'{{{a_ns}}}bodyPr')
                if bodyPr is not None:
                    bodyPr.set('anchor', 'ctr')

    log.info("Theme-Text-Layout angepasst")


def _build_slide_plan(song_paths: dict, skip_slides: set = None,
                      section_order: list = None,
                      extra_song_paths: dict = None) -> list:
    """Erstellt den Folienplan: welche Folien in welcher Reihenfolge.

    Template hat 37 Folien (0-36). Abschnitte werden in der gegebenen
    Reihenfolge (section_order) angeordnet. An Song-Positionen werden
    Platzhalter durch Lied-Dateien ersetzt.

    Args:
        song_paths: Dict {slot_key: pfad} für die regulären Lieder
        skip_slides: Set von 0-basierten Indizes die übersprungen werden
        section_order: Liste von Section-Keys in gewünschter Reihenfolge
        extra_song_paths: Dict {song_extra1: pfad, ...} für Extra-Lieder
    """
    plan = []
    song_slot_indices = set(SONG_SLOTS.values())
    skip = skip_slides or set()
    order = section_order or DEFAULT_SECTION_ORDER

    # Disabled sections: alle Slides dieser Sections zu skip hinzufügen
    disabled_sections = set()
    for section_key in TOGGLEABLE_SECTIONS:
        section_slides = set(TOGGLEABLE_SECTIONS[section_key]["slides"])
        if section_slides.issubset(skip):
            disabled_sections.add(section_key)

    def _add_slide(i):
        """Fügt eine einzelne Folie zum Plan hinzu (mit Song-Ersetzung)."""
        if i in skip:
            log.info(f"Plan: Folie {i+1} übersprungen (deaktiviert)")
            return
        if i in song_slot_indices:
            slot_key = [k for k, v in SONG_SLOTS.items() if v == i][0]
            song_path = song_paths.get(slot_key)
            if song_path and os.path.exists(song_path):
                plan.append(("file", song_path))
                log.info(f"Plan: {slot_key} → {os.path.basename(song_path)}")
            else:
                plan.append(("template", i))
                log.warning(f"Plan: {slot_key} → Platzhalter (kein Lied)")
        else:
            plan.append(("template", i))

    # 1. Intro-Folien (immer am Anfang)
    for i in INTRO_SLIDES:
        _add_slide(i)

    # 2. Abschnitte in der gewünschten Reihenfolge
    extras = extra_song_paths or {}
    extras_in_order = set()  # Track welche Extras explizit platziert wurden

    for section_key in order:
        # Extra-Lied an dieser Position einfügen
        if section_key.startswith("song_extra"):
            extras_in_order.add(section_key)
            if section_key in disabled_sections:
                log.info(f"Plan: {section_key} deaktiviert, übersprungen")
                continue
            extra_path = extras.get(section_key)
            if extra_path and os.path.exists(extra_path):
                # Thema-Folie + Extra-Lied + Thema-Folie
                plan.append(("template", INTRO_SLIDES[-1]))
                plan.append(("file", extra_path))
                plan.append(("template", INTRO_SLIDES[-1]))
                log.info(f"Plan: {section_key} → {os.path.basename(extra_path)} (Extra)")
            else:
                log.warning(f"Plan: {section_key} → kein Pfad gefunden")
            continue

        if section_key not in SECTION_BLOCKS:
            log.warning(f"Plan: Unbekannter Abschnitt '{section_key}', übersprungen")
            continue
        if section_key in disabled_sections:
            log.info(f"Plan: Abschnitt '{section_key}' deaktiviert, übersprungen")
            continue

        block_slides = SECTION_BLOCKS[section_key]
        for i in block_slides:
            _add_slide(i)

    # Extras die nicht in der Order waren, am Ende einfügen
    for extra_key in sorted(extras.keys()):
        if extra_key not in extras_in_order:
            extra_path = extras[extra_key]
            if extra_path and os.path.exists(extra_path):
                plan.append(("template", INTRO_SLIDES[-1]))
                plan.append(("file", extra_path))
                plan.append(("template", INTRO_SLIDES[-1]))
                log.info(f"Plan: {extra_key} → {os.path.basename(extra_path)} (Extra, Ende)")

    return plan


def _fill_all_text(prs, data, fetch_bible: bool, template_indices: set = None):
    """Füllt alle Textplatzhalter in der Präsentation.

    Da die Folienreihenfolge jetzt dynamisch ist (durch eingefügte Lieder),
    suchen wir Folien anhand ihres Textinhalts statt nach fester Position.
    Bibeltexte werden bei Bedarf auf mehrere Folien aufgeteilt.
    """
    # Pre-fetch bible texts so we know how many slides we need
    bible_data = {}
    if fetch_bible:
        if data.lesung_reference:
            bible_data['lesung'] = fetch_bible_text(data.lesung_reference)
        if data.predigt1_reference:
            bible_data['predigt1'] = fetch_bible_text(data.predigt1_reference)
        if data.predigt2_reference:
            bible_data['predigt2'] = fetch_bible_text(data.predigt2_reference)

    # First pass: find bible text slides and duplicate if needed
    # We iterate by index because we insert slides during the loop
    i = 0
    while i < len(prs.slides):
        slide = prs.slides[i]
        slide_text = _get_slide_text(slide)

        bible_key = None
        placeholder_text = None
        header_old = None
        header_new = None

        if "Lesung: Bibelstelle" in slide_text:
            bible_key = 'lesung'
            placeholder_text = "Lesungstext"
            header_old = "Lesung: Bibelstelle"
            header_new = f"Lesung: {data.lesung_reference}"
        elif "Predigttext1" in slide_text and "Predigt: Bibelstelle" in slide_text:
            bible_key = 'predigt1'
            placeholder_text = "Predigttext1"
            header_old = "Predigt: Bibelstelle"
            header_new = f"Predigt: {data.predigt1_reference}"
        elif "Predigttext2" in slide_text and "Predigt: Bibelstelle" in slide_text:
            bible_key = 'predigt2'
            placeholder_text = "Predigttext2"
            header_old = "Predigt: Bibelstelle"
            header_new = f"Predigt : {data.predigt2_reference}"

        if bible_key and bible_key in bible_data:
            verses = bible_data[bible_key]
            chunks = _split_verses_for_slides(verses)
            total_pages = len(chunks)

            # ERST duplizieren (solange Folie noch Platzhalter hat)
            extra_slides = []
            for _ in range(len(chunks) - 1):
                new_slide = _duplicate_slide_after(
                    prs, i + len(extra_slides), slide)
                extra_slides.append(new_slide)

            # Template-Indices verschieben (neue Folien eingefügt)
            if template_indices and extra_slides:
                n_extra = len(extra_slides)
                insert_pos = i + 1
                shifted = set()
                for idx in template_indices:
                    if idx >= insert_pos:
                        shifted.add(idx + n_extra)
                    else:
                        shifted.add(idx)
                template_indices.clear()
                template_indices.update(shifted)

            # DANN alle Folien füllen (Original + Duplikate)
            all_bible_slides = [slide] + extra_slides
            for ci, (bible_slide, chunk) in enumerate(
                    zip(all_bible_slides, chunks), start=1):
                _replace_text_in_slide(bible_slide, header_old, header_new)
                _insert_bible_verses(bible_slide, placeholder_text, chunk)
                # Seitenzahl aktualisieren (Template hat "1 / 3" oder "1 / 2")
                _update_page_number(bible_slide, ci, total_pages)

            i += len(extra_slides)  # Skip duplicated slides

        else:
            # Non-bible slides: fill text normally
            is_template = template_indices and i in template_indices
            _fill_slide_simple(slide, data, is_template)

        i += 1


def _fix_title_overflow(shape):
    """Verhindert, dass der Titel über den Folienrand hinausragt.

    Vergrößert die Text-Box ins obere Drittel der Folie (0.5 - 6.5cm)
    und aktiviert Word-Wrap. Anker auf Top, damit Text von oben nach
    unten wächst und nie über den oberen Rand ragt.
    """
    try:
        from pptx.util import Cm

        # Oberes Drittel: 0.5cm bis 6.5cm (6cm Höhe, genug für 2 Zeilen à 66pt)
        shape.top = Cm(0.5)
        shape.height = Cm(6)

        # Word-Wrap aktivieren damit Text umbricht
        shape.text_frame.word_wrap = True

        # Anker auf Top setzen: Text beginnt oben und wächst nach unten
        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 't')

    except Exception as e:
        log.debug(f"_fix_title_overflow fehlgeschlagen: {e}")


def _build_standard_rPr(sz, lang='de-DE', baseline=None):
    """Erzeugt ein standardisiertes rPr-Element (identisch für alle Folien).

    Args:
        sz: Schriftgröße in Hunderstel-Punkt (z.B. '6600' für 66pt)
        lang: Sprache
        baseline: Hochstellung (z.B. '30000' für Versnummern)
    """
    rPr_xml = (
        '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' kumimoji="0" lang="{lang}" sz="{sz}"'
        ' b="0" i="0" u="none" strike="noStrike"'
        ' kern="1200" cap="none" spc="0" normalizeH="0"'
        ' baseline="{baseline}" noProof="0" dirty="0">'
        '<a:ln><a:noFill/></a:ln>'
        '<a:solidFill><a:prstClr val="white"/></a:solidFill>'
        '<a:effectLst>'
        '<a:outerShdw blurRad="50800" dist="38100" dir="2700000"'
        ' algn="tl" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="70000"/></a:srgbClr>'
        '</a:outerShdw>'
        '</a:effectLst>'
        '<a:uLnTx/>'
        '<a:uFillTx/>'
        '<a:latin typeface="Baskerville Old Face"'
        ' panose="02020602080505020303" pitchFamily="18" charset="0"/>'
        '<a:ea typeface="+mj-ea"/>'
        '<a:cs typeface="+mj-cs"/>'
        '</a:rPr>'
    ).format(lang=lang, sz=sz, baseline=baseline or '0')
    return etree.fromstring(rPr_xml)


def _has_white_text(shape):
    """Prüft ob ein Shape weißen Text hat (prstClr=white oder srgbClr=FFFFFF)."""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                continue
            fill = rPr.find(qn('a:solidFill'))
            if fill is not None:
                prstClr = fill.find(qn('a:prstClr'))
                if prstClr is not None and prstClr.get('val') == 'white':
                    return True
                srgbClr = fill.find(qn('a:srgbClr'))
                if srgbClr is not None and srgbClr.get('val', '').upper() == 'FFFFFF':
                    return True
    return False


def _add_text_shadow(slide):
    """Normalisiert Formatierung auf Runs mit weißem Text (Schatten, Schrift, Stil).

    Shapes mit nicht-weißem Text (z.B. Glaubensbekenntnis) bleiben unverändert.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not _has_white_text(shape):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                old_rPr = run._r.find(qn('a:rPr'))
                # Schriftgröße und Baseline beibehalten
                sz = old_rPr.get('sz', '5400') if old_rPr is not None else '5400'
                lang = old_rPr.get('lang', 'de-DE') if old_rPr is not None else 'de-DE'
                baseline = old_rPr.get('baseline') if old_rPr is not None else None

                # Altes rPr entfernen und durch Standard ersetzen
                if old_rPr is not None:
                    run._r.remove(old_rPr)
                new_rPr = _build_standard_rPr(sz, lang, baseline)
                run._r.insert(0, new_rPr)


def _set_cell_text(tc, text):
    """Setzt den Text einer Tabellenzelle (a:tc Element), behält Formatierung."""
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    txBody = tc.find(qn('a:txBody'))
    if txBody is None:
        return
    para = txBody.find(qn('a:p'))
    if para is None:
        return
    run = para.find(qn('a:r'))
    if run is not None:
        t = run.find(qn('a:t'))
        if t is not None:
            t.text = text
    else:
        # Kein Run vorhanden → neuen erstellen
        r = etree.SubElement(para, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', 'de-DE')
        rPr.set('sz', '2800')
        rPr.set('b', '0')
        t = etree.SubElement(r, qn('a:t'))
        t.text = text


def _fill_invitation_slide(slide, events):
    """Befüllt die Tabelle auf der 'Herzliche Einladung' Folie mit Events.

    Args:
        slide: Die Folie mit der Einladungstabelle
        events: Liste von InvitationEvent-Objekten
    """
    # Tabelle finden
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table_shape = shape
            break

    if table_shape is None:
        log.warning("Keine Tabelle auf Herzliche Einladung Folie gefunden")
        return

    table = table_shape.table
    tbl = table._tbl
    rows = tbl.findall(qn('a:tr'))

    if len(rows) < 2:
        log.warning("Tabelle hat weniger als 2 Zeilen")
        return

    # Datum-Spalte verbreitern damit "Di 17.03.26" in eine Zeile passt
    # Col 0: 4.9cm → 6.5cm, Col 2: 15.0cm → 13.4cm (Differenz umverteilen)
    from pptx.util import Cm
    table.columns[0].width = Cm(6.5)
    table.columns[2].width = Cm(13.4)

    # Template-Datenzeile klonen (Row 1) für Formatierung
    row_template = copy.deepcopy(rows[1])

    # Alle Datenzeilen entfernen (Header = rows[0] bleibt)
    for row in rows[1:]:
        tbl.remove(row)

    # Neue Zeilen einfügen
    for event in events:
        new_row = copy.deepcopy(row_template)
        cells = new_row.findall(qn('a:tc'))
        if len(cells) >= 3:
            _set_cell_text(cells[0], event.date_str)
            _set_cell_text(cells[1], event.time_str)
            # Hinweis in gleiche Zeile anhängen
            veranstaltung = event.event_name
            if event.note:
                veranstaltung += f" \u2014 {event.note}"
            _set_cell_text(cells[2], veranstaltung)
        tbl.append(new_row)

    # "----Optional---" Textbox entfernen
    for shape in slide.shapes:
        if shape.has_text_frame and "----Optional---" in shape.text_frame.text:
            sp = shape._element
            sp.getparent().remove(sp)
            break

    log.info(f"Herzliche Einladung: {len(events)} Events in Tabelle eingefügt")


def _adjust_song_text_size(slide, min_sz=3200):
    """Hebt zu kleine Liedtexte auf mindestens min_sz (32pt) an.

    Betrifft nur Liedtext-Runs (sz >= 2000), nicht Header, Copyright
    oder Seitenzahlen (die haben sz < 2000).
    """
    min_str = str(min_sz)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    continue
                sz = rPr.get('sz')
                if sz and sz.isdigit() and 2000 <= int(sz) < min_sz:
                    rPr.set('sz', min_str)


def _fill_slide_simple(slide, data, is_template=True):
    """Füllt einfache Text-Platzhalter (nicht Bibeltexte)."""
    slide_text = _get_slide_text(slide)

    if "Herzliche Einladung" in slide_text:
        _fill_invitation_slide(slide, data.invitation_events)
        return

    if "Gottes Liebe erkennen" in slide_text:
        _replace_text_in_slide(slide, "Gottes Liebe erkennen", data.theme)
        _replace_text_in_slide(slide, "Gottesdienst am 26.11.2023,",
                               f"Gottesdienst am {data.date_str},")
        _replace_text_in_slide(slide, "Letzter Sonntag des Kirchenjahres",
                               data.kirchenkalender)
        # Titel-Box vergrößern damit kein Overflow über Folienrand
        for shape in slide.shapes:
            if shape.has_text_frame and data.theme in shape.text_frame.text:
                _fix_title_overflow(shape)

    if "Begrüßungsvers" in slide_text and data.greeting_verse:
        _replace_text_in_slide(slide, "Begrüßungsvers", data.greeting_verse)

    if "Gottesdienst am XY.XY.XY" in slide_text:
        _replace_text_in_slide(slide, "Thema", data.theme)
        _replace_text_in_slide(slide, "Gottesdienst am XY.XY.XY",
                               f"Gottesdienst am {data.date_str},")
        _replace_text_in_slide(slide, "Name des Sonntags", data.kirchenkalender)
        # Titel-Box vergrößern damit kein Overflow über Folienrand
        for shape in slide.shapes:
            if shape.has_text_frame and data.theme in shape.text_frame.text:
                _fix_title_overflow(shape)

    if "Predigt zu Predigtext1" in slide_text:
        _replace_text_in_slide(slide, "Predigt zu Predigtext1",
                               f"Predigt zu {data.predigt1_reference}")
        _replace_text_in_slide(slide, "Thema Predigt1", data.predigt1_title)

    if "Predigt zu Predigttext2" in slide_text:
        _replace_text_in_slide(slide, "Predigt zu Predigttext2",
                               f"Predigt zu {data.predigt2_reference}")
        _replace_text_in_slide(slide, "Predigtthema 2", data.predigt2_title)

    # Text-Schatten nur auf Template-Folien (nicht auf reinkopierte Lieder etc.)
    if is_template:
        _add_text_shadow(slide)

    # Lied-Folien: Liedtext auf 32pt anpassen (Header bleibt unverändert)
    if not is_template:
        _adjust_song_text_size(slide)


def _update_page_number(slide, current_page, total_pages):
    """Aktualisiert die Seitenzahl-Anzeige (z.B. '1 / 3' → '2 / 2')."""
    import re
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if re.match(r'^\d+\s*/\s*\d+$', text):
                if para.runs:
                    para.runs[0].text = f"{current_page} / {total_pages}"
                    for run in para.runs[1:]:
                        run.text = ""
                return


def _split_verses_for_slides(verses, max_chars=350):
    """Teilt Verse in Gruppen auf, die jeweils auf eine Folie passen.

    Args:
        verses: Liste von (vers_nummer, vers_text) Tupeln
        max_chars: Maximale Zeichenanzahl pro Folie (bei 32pt Cambria)

    Returns:
        Liste von Listen von (vers_nummer, vers_text) Tupeln
    """
    if not verses:
        return [verses]

    chunks = []
    current_chunk = []
    current_len = 0

    for v_num, v_text in verses:
        verse_len = len(str(v_num)) + len(v_text) + 1
        if current_len + verse_len > max_chars and current_chunk:
            chunks.append(current_chunk)
            current_chunk = []
            current_len = 0
        current_chunk.append((v_num, v_text))
        current_len += verse_len

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def _duplicate_slide_after(prs, slide_index, source_slide):
    """Dupliziert eine Folie und fügt die Kopie direkt danach ein.

    Kopiert die Original-Folie (mit Layout, Shapes, Bildern) und
    verschiebt sie an die richtige Position.
    """
    # Kopiere die Folie ans Ende (copy_slide fügt immer am Ende an)
    new_slide = copy_slide(prs, source_slide, prs, is_song=False)

    # Verschiebe von Ende an die Position nach slide_index
    _move_slide_to(prs, len(prs.slides) - 1, slide_index + 1)

    return new_slide


def _move_slide_to(prs, old_index, new_index):
    """Verschiebt eine Folie von old_index nach new_index."""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[old_index]
    sldIdLst.remove(sldId)
    if new_index >= len(sldIdLst):
        sldIdLst.append(sldId)
    else:
        sldIdLst.insert(new_index, sldId)


def _build_bible_rPr(baseline=None):
    """Erzeugt ein rPr für Bibeltext: Cambria, 32pt, nicht fett, schwarz."""
    rPr_xml = (
        '<a:rPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' lang="de-DE" sz="3200" b="0" i="0" dirty="0"'
        ' baseline="{baseline}">'
        '<a:solidFill><a:srgbClr val="000000"/></a:solidFill>'
        '<a:latin typeface="Cambria"/>'
        '</a:rPr>'
    ).format(baseline=baseline or '0')
    return etree.fromstring(rPr_xml)


def _insert_bible_verses(slide, placeholder_text, verses):
    """Ersetzt einen Platzhalter-Text durch Bibelverse mit hochgestellten Nummern.

    Verwendet Cambria 32pt, nicht fett, weiß mit Schatten.

    Args:
        slide: Ziel-Folie
        placeholder_text: Text der ersetzt werden soll (z.B. "Lesungstext")
        verses: Liste von (vers_nummer, vers_text) Tupeln
    """
    from pptx.oxml.ns import qn

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if placeholder_text not in para.text:
                continue

            # Alle bestehenden Runs entfernen
            for run in list(para.runs):
                run._r.getparent().remove(run._r)

            # endParaRPr temporär entfernen (muss am Ende stehen)
            endParaRPr = para._p.find(qn('a:endParaRPr'))
            if endParaRPr is not None:
                para._p.remove(endParaRPr)

            # Verse einfügen mit hochgestellten Nummern
            for v_num, v_text in verses:
                if v_num > 0:
                    # Vers-Nummer hochgestellt
                    num_r = etree.SubElement(para._p, qn('a:r'))
                    num_r.insert(0, _build_bible_rPr(baseline='30000'))
                    num_t = etree.SubElement(num_r, qn('a:t'))
                    num_t.text = str(v_num)

                # Vers-Text
                text_r = etree.SubElement(para._p, qn('a:r'))
                text_r.insert(0, _build_bible_rPr())
                text_t = etree.SubElement(text_r, qn('a:t'))
                text_t.text = v_text + " "

            # endParaRPr wieder ans Ende setzen
            if endParaRPr is not None:
                para._p.append(endParaRPr)

            log.info(f"Bibeltext eingefügt: {len(verses)} Verse für '{placeholder_text}'")
            return


def _get_slide_text(slide) -> str:
    """Gibt den gesamten Text einer Folie als String zurück."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    from datetime import date
    from excel_reader import read_godi_plan
    from song_finder import build_song_index, find_song

    sunday = date(2026, 3, 8)
    data = read_godi_plan(sunday)
    if not data:
        print("GoDi-Plan nicht gefunden!")
        exit(1)

    index = build_song_index()
    song_paths = {}
    for song in data.songs:
        path = find_song(song, index)
        if path:
            song_paths[song.slot_key] = path

    output = build_presentation(data, song_paths, fetch_bible=True)
    print(f"\nPräsentation erstellt: {output}")

    # Verify
    prs = Presentation(output)
    print(f"Folien: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t[:65])
                        break
        print(f"  {i+1:3d}: {texts[:2]}")
